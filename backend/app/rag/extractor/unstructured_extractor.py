"""Unstructured Office 解析器

使用 python-docx、python-pptx 等库解析 Office 文件：
- .doc, .docx (Word) - 使用 python-docx
- .ppt, .pptx (PowerPoint) - 使用 python-pptx
- .xls, .xlsx (Excel) - 使用 openpyxl
- .html (HTML) - 使用 BeautifulSoup
"""

from typing import List
from pathlib import Path

from app.rag.extractor.base_extractor import BaseExtractor, ContentBlock


class UnstructuredExtractor(BaseExtractor):
    """Unstructured Office 解析器"""

    SUPPORTED_EXTENSIONS = [
        ".doc", ".docx",
        ".ppt", ".pptx",
        ".xls", ".xlsx",
        ".html", ".htm",
    ]

    async def extract(self, file_path: str) -> List[ContentBlock]:
        """使用 python-docx 等库解析 Office 文件

        Args:
            file_path: 文件路径

        Returns:
            内容块列表
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        blocks = []

        # Word 文档 (.docx)
        if ext in [".doc", ".docx"]:
            blocks = await self._extract_docx(file_path)

        # PowerPoint (.pptx)
        elif ext in [".ppt", ".pptx"]:
            blocks = await self._extract_pptx(file_path)

        # Excel (.xlsx)
        elif ext in [".xls", ".xlsx"]:
            blocks = await self._extract_xlsx(file_path)

        # HTML
        elif ext in [".html", ".htm"]:
            blocks = await self._extract_html(file_path)

        return blocks

    async def _extract_docx(self, file_path: str) -> List[ContentBlock]:
        """解析 Word 文档

        Args:
            file_path: 文件路径

        Returns:
            内容块列表
        """
        try:
            from docx import Document
        except ImportError:
            # python-docx 未安装，返回提示
            return [
                ContentBlock(
                    type="text",
                    content="Word 文档解析失败：需要安装 python-docx 库\n请运行: pip install python-docx",
                    page_number=1,
                    metadata={"source": Path(file_path).name, "error": "missing_library"}
                )
            ]

        try:
            # 方式1: 标准解析
            doc = Document(file_path)
            blocks = []
            page_number = 1

            # 提取所有段落
            for i, para in enumerate(doc.paragraphs):
                text = para.text.strip()
                if text:
                    # 检测段落类型（标题、正文等）
                    para_type = self._detect_paragraph_type(para)

                    blocks.append(ContentBlock(
                        type="text",
                        content=text,
                        page_number=page_number,
                        metadata={
                            "paragraph_index": i,
                            "style": para.style.name if para.style else "Normal",
                            "para_type": para_type,
                            "source": Path(file_path).name,
                            "extractor": "python-docx"
                        }
                    ))

            # 提取表格
            for i, table in enumerate(doc.tables):
                # 将表格转换为 Markdown 格式
                table_md = self._table_to_markdown(table)

                blocks.append(ContentBlock(
                    type="table",
                    content=table_md,
                    page_number=page_number,
                    metadata={
                        "table_index": i,
                        "rows": len(table.rows),
                        "cols": len(table.columns),
                        "source": Path(file_path).name,
                        "extractor": "python-docx"
                    }
                ))

            # 如果成功提取到内容，返回结果
            if blocks:
                return blocks

            # 如果标准解析没有提取到内容，尝试备用方法
            return await self._extract_docx_fallback(file_path)

        except Exception as e:
            # 标准解析失败，尝试备用方法
            import logging
            logger = logging.getLogger(__name__)
            logger.warning(f"Standard docx parsing failed: {e}. Trying fallback method.")

            return await self._extract_docx_fallback(file_path)

    async def _extract_pptx(self, file_path: str) -> List[ContentBlock]:
        """解析 PowerPoint 文档

        Args:
            file_path: 文件路径

        Returns:
            内容块列表
        """
        try:
            from pptx import Presentation
        except ImportError:
            return [
                ContentBlock(
                    type="text",
                    content="PowerPoint 解析失败：需要安装 python-pptx 库\n请运行: pip install python-pptx",
                    page_number=1,
                    metadata={"source": Path(file_path).name, "error": "missing_library"}
                )
            ]

        try:
            prs = Presentation(file_path)
            blocks = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_content = []

                # 提取所有形状中的文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_content.append(text)

                if slide_content:
                    blocks.append(ContentBlock(
                        type="text",
                        content="\n\n".join(slide_content),
                        page_number=slide_num,
                        metadata={
                            "slide_number": slide_num,
                            "source": Path(file_path).name,
                            "extractor": "python-pptx"
                        }
                    ))

            return blocks

        except Exception as e:
            return [
                ContentBlock(
                    type="text",
                    content=f"PowerPoint 解析失败：{str(e)}",
                    page_number=1,
                    metadata={"source": Path(file_path).name, "error": str(e)}
                )
            ]

    async def _extract_xlsx(self, file_path: str) -> List[ContentBlock]:
        """解析 Excel 文档

        Args:
            file_path: 文件路径

        Returns:
            内容块列表
        """
        try:
            from openpyxl import load_workbook
        except ImportError:
            return [
                ContentBlock(
                    type="text",
                    content="Excel 解析失败：需要安装 openpyxl 库\n请运行: pip install openpyxl",
                    page_number=1,
                    metadata={"source": Path(file_path).name, "error": "missing_library"}
                )
            ]

        try:
            wb = load_workbook(file_path, data_only=True)
            blocks = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]

                # 将sheet转换为Markdown表格
                rows_data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        rows_data.append(row)

                if rows_data:
                    # 构建Markdown表格
                    table_md = self._rows_to_markdown_table(rows_data)

                    blocks.append(ContentBlock(
                        type="table",
                        content=table_md,
                        page_number=1,
                        metadata={
                            "sheet_name": sheet_name,
                            "rows": len(rows_data),
                            "source": Path(file_path).name,
                            "extractor": "openpyxl"
                        }
                    ))

            return blocks

        except Exception as e:
            return [
                ContentBlock(
                    type="text",
                    content=f"Excel 解析失败：{str(e)}",
                    page_number=1,
                    metadata={"source": Path(file_path).name, "error": str(e)}
                )
            ]

    async def _extract_html(self, file_path: str) -> List[ContentBlock]:
        """解析 HTML 文档

        Args:
            file_path: 文件路径

        Returns:
            内容块列表
        """
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            return [
                ContentBlock(
                    type="text",
                    content="HTML 解析失败：需要安装 beautifulsoup4 库\n请运行: pip install beautifulsoup4",
                    page_number=1,
                    metadata={"source": Path(file_path).name, "error": "missing_library"}
                )
            ]

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, 'html.parser')

            # 提取正文内容
            # 移除script和style标签
            for script in soup(["script", "style"]):
                script.decompose()

            text = soup.get_text(separator='\n', strip=True)

            # 按段落分割
            paragraphs = [p.strip() for p in text.split('\n') if p.strip()]

            blocks = []
            for i, para in enumerate(paragraphs):
                blocks.append(ContentBlock(
                    type="text",
                    content=para,
                    page_number=1,
                    metadata={
                        "paragraph_index": i,
                        "source": Path(file_path).name,
                        "extractor": "beautifulsoup4"
                    }
                ))

            return blocks

        except Exception as e:
            return [
                ContentBlock(
                    type="text",
                    content=f"HTML 解析失败：{str(e)}",
                    page_number=1,
                    metadata={"source": Path(file_path).name, "error": str(e)}
                )
            ]

    def _detect_paragraph_type(self, para) -> str:
        """检测段落类型

        Args:
            para: docx Paragraph 对象

        Returns:
            段落类型: title, heading, list, text
        """
        style_name = para.style.name if para.style else "Normal"

        if "Title" in style_name or "Heading 1" in style_name:
            return "title"
        elif "Heading" in style_name:
            return "heading"
        elif "List" in style_name:
            return "list"
        else:
            return "text"

    def _table_to_markdown(self, table) -> str:
        """将 Word 表格转换为 Markdown 格式

        Args:
            table: docx Table 对象

        Returns:
            Markdown 表格字符串
        """
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        # 构建Markdown表格
        md_lines = []

        # 表头
        header = rows[0]
        md_lines.append("| " + " | ".join(header) + " |")

        # 分隔线
        md_lines.append("| " + " | ".join(["---"] * len(header)) + " |")

        # 数据行
        for row in rows[1:]:
            md_lines.append("| " + " | ".join(row) + " |")

        return "\n".join(md_lines)

    def _rows_to_markdown_table(self, rows) -> str:
        """将数据行转换为 Markdown 表格

        Args:
            rows: 数据行列表

        Returns:
            Markdown 表格字符串
        """
        if not rows:
            return ""

        # 转换None为空字符串
        rows_str = [[str(cell) if cell is not None else "" for cell in row] for row in rows]

        md_lines = []

        # 表头
        header = rows_str[0]
        md_lines.append("| " + " | ".join(header) + " |")

        # 分隔线
        md_lines.append("| " + " | ".join(["---"] * len(header)) + " |")

        # 数据行
        for row in rows_str[1:]:
            md_lines.append("| " + " | ".join(row) + " |")

        return "\n".join(md_lines)

    @classmethod
    def supports_file_type(cls, file_extension: str) -> bool:
        """检查是否支持"""
        return file_extension.lower() in cls.SUPPORTED_EXTENSIONS

    async def _extract_docx_fallback(self, file_path: str) -> List[ContentBlock]:
        """备用docx解析方法（使用zipfile直接解析）

        当python-docx失败时，尝试直接解析XML

        Args:
            file_path: 文件路径

        Returns:
            内容块列表
        """
        import zipfile
        import xml.etree.ElementTree as ET

        try:
            # DOCX实际上是一个ZIP文件，包含XML
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                # 读取document.xml
                if 'word/document.xml' in zip_ref.namelist():
                    with zip_ref.open('word/document.xml') as xml_file:
                        tree = ET.parse(xml_file)
                        root = tree.getroot()

                        # 提取所有文本节点
                        # Word XML namespace
                        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

                        texts = []
                        for elem in root.iter():
                            # 查找所有<w:t>标签
                            if elem.tag == '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t':
                                if elem.text:
                                    texts.append(elem.text)

                        # 合并文本
                        full_text = '\n'.join(texts)

                        if full_text.strip():
                            return [
                                ContentBlock(
                                    type="text",
                                    content=full_text,
                                    page_number=1,
                                    metadata={
                                        "source": Path(file_path).name,
                                        "extractor": "zipfile_xml_fallback",
                                        "method": "direct_xml_parse"
                                    }
                                )
                            ]

        except Exception as fallback_error:
            # 备用方法也失败
            return [
                ContentBlock(
                    type="text",
                    content=f"Word 文档解析失败（所有方法都失败）\n错误1: {str(e)}\n错误2: {str(fallback_error)}",
                    page_number=1,
                    metadata={
                        "source": Path(file_path).name,
                        "error": f"{str(e)}; {str(fallback_error)}",
                        "extractor": "failed"
                    }
                )
            ]

        # 如果没有提取到任何内容
        return [
            ContentBlock(
                type="text",
                content=f"Word 文档解析失败：未能提取到任何内容",
                page_number=1,
                metadata={
                    "source": Path(file_path).name,
                    "error": "no_content_extracted",
                    "extractor": "failed"
                }
            )
        ]